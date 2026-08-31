import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".pptdeps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SOURCE = r"C:\Users\Ian\Documents\4th yr output\studymate AI presentation.pptx"
OUTPUT = os.path.join(os.path.dirname(__file__), "studymate AI presentation - updated.pptx")

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(224, 242, 254)
TEAL = RGBColor(13, 148, 136)
WHITE = RGBColor(255, 255, 255)
SLATE = RGBColor(71, 85, 105)
LIGHT = RGBColor(248, 250, 252)
LINE = RGBColor(203, 213, 225)


def rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def text(slide, value, x, y, w, h, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def base_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 13.333, 7.5, WHITE)
    rect(slide, 0, 0, 13.333, 0.18, TEAL)
    text(slide, title, 0.7, 0.45, 11.9, 0.5, 28, NAVY, True)
    text(slide, subtitle, 0.72, 0.99, 11.8, 0.38, 12, SLATE)
    text(slide, "StudyMate AI • Implemented web prototype", 0.72, 7.12, 6, 0.2, 9, SLATE)
    return slide


def bullet_card(slide, heading, points, x, y, w, h, accent=BLUE):
    rect(slide, x, y, w, h, LIGHT, True, LINE)
    rect(slide, x, y, 0.09, h, accent, False)
    text(slide, heading, x + 0.28, y + 0.2, w - 0.5, 0.32, 16, NAVY, True)
    top = y + 0.7
    for point in points:
        text(slide, "• " + point, x + 0.3, top, w - 0.55, 0.42, 12.5, SLATE)
        top += 0.52


def phone_mockup(slide, x, y, w, h):
    rect(slide, x, y, w, h, NAVY, True)
    rect(slide, x + 0.15, y + 0.18, w - 0.3, h - 0.36, WHITE, True)
    text(slide, "StudyMate", x + 0.35, y + 0.4, w - 0.7, 0.26, 14, NAVY, True, PP_ALIGN.CENTER)
    rect(slide, x + 0.35, y + 0.92, w - 0.7, 0.58, SKY, True)
    text(slide, "Hello! How can I help you today?", x + 0.48, y + 1.05, w - 0.95, 0.27, 8.5, NAVY)
    rect(slide, x + 0.62, y + 1.75, w - 1.0, 0.54, RGBColor(219, 234, 254), True)
    text(slide, "Summarize my notes", x + 0.75, y + 1.87, w - 1.25, 0.22, 8.5, NAVY)
    rect(slide, x + 0.35, y + h - 0.85, w - 0.7, 0.42, LIGHT, True, LINE)
    text(slide, "Attach  •  Ask StudyMate", x + 0.47, y + h - 0.76, w - 0.95, 0.17, 7.5, SLATE)


def main():
    prs = Presentation(SOURCE)

    slide = base_slide(prs, "Implemented Prototype", "The current StudyMate AI web application extends the proposed platform with usable student-facing tools.")
    phone_mockup(slide, 0.9, 1.65, 3.1, 4.9)
    bullet_card(slide, "What is working in the prototype", [
        "AI chat interface with saved conversation history",
        "Document attachment and text extraction before prompting",
        "Voice input and read-aloud response controls",
        "Projects, study schedule, downloads, and photo editing",
    ], 4.55, 1.65, 7.75, 2.75, TEAL)
    bullet_card(slide, "Student-centered experience", [
        "Responsive side navigation keeps learning tools organized.",
        "Browser storage retains chats, projects, and scheduled sessions.",
    ], 4.55, 4.72, 7.75, 1.55, BLUE)

    slide = base_slide(prs, "AI Chat and Document Support", "A focused workspace where students can ask questions, attach study material, and review replies.")
    bullet_card(slide, "Chat interaction", [
        "Send questions through the StudyMate chat panel.",
        "Start a new conversation or reopen a saved chat from History.",
        "Optional text-to-speech reads AI responses aloud.",
        "Speech recognition can turn a spoken question into a prompt.",
    ], 0.75, 1.6, 5.85, 4.75, BLUE)
    bullet_card(slide, "Supported study files", [
        "PDF and DOCX documents",
        "TXT, Markdown, CSV, JSON, and common source-code files",
        "Extracted text is included with the learner’s question for context.",
    ], 6.95, 1.6, 5.6, 4.75, TEAL)

    slide = base_slide(prs, "Learning Organization and Exports", "Practical features help learners retain work and plan their next study session.")
    bullet_card(slide, "Organize", [
        "Create named study projects.",
        "Add upcoming sessions with a topic and date/time.",
        "Review and delete saved conversations from History.",
    ], 0.75, 1.55, 3.9, 4.65, TEAL)
    bullet_card(slide, "Download", [
        "Export the active conversation as TXT.",
        "Generate a shareable PDF copy.",
        "Create a DOCX version for notes and review.",
    ], 4.72, 1.55, 3.9, 4.65, BLUE)
    bullet_card(slide, "Photo Studio", [
        "Upload PNG, JPEG, or WebP study images.",
        "Adjust brightness, contrast, saturation, and grayscale.",
        "Download the edited image as PNG or JPG.",
    ], 8.69, 1.55, 3.9, 4.65, RGBColor(124, 58, 237))

    slide = base_slide(prs, "Current Prototype Technology", "The implemented version is a lightweight client-side web application; these details reflect the source files in the project folder.")
    bullet_card(slide, "Interface", [
        "HTML5 structure and responsive CSS styling",
        "Modern ES module JavaScript for application behavior",
        "Canvas-based image adjustments and browser downloads",
    ], 0.75, 1.55, 3.9, 4.8, BLUE)
    bullet_card(slide, "Browser services", [
        "Local storage for chats, projects, and study schedules",
        "Web Speech APIs for recognition and text-to-speech",
        "Client-side readers for PDF and DOCX attachments",
    ], 4.72, 1.55, 3.9, 4.8, TEAL)
    bullet_card(slide, "AI connection", [
        "Google Gemini generate-content request",
        "Attached file text is supplied as prompt context",
        "API key configuration is kept separate from UI logic",
    ], 8.69, 1.55, 3.9, 4.8, RGBColor(124, 58, 237))

    slide = base_slide(prs, "Prototype Alignment and Next Steps", "The implemented interface provides a concrete foundation for the broader architecture proposed in this study.")
    bullet_card(slide, "Already demonstrated", [
        "Conversational academic support workflow",
        "File-assisted study prompts and accessible voice controls",
        "Personal study organization and exported conversation notes",
    ], 0.75, 1.6, 5.85, 4.65, TEAL)
    bullet_card(slide, "Recommended next development phase", [
        "Move the API key to a secure server-side service.",
        "Add authenticated user accounts and a persistent database.",
        "Expand document analysis into summaries, quizzes, and flashcards.",
    ], 6.95, 1.6, 5.6, 4.65, BLUE)

    prs.save(OUTPUT)
    print(OUTPUT)
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
